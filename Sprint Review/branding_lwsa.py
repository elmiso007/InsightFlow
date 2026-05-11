"""
Extrai paleta e imagens do template PowerPoint LWSA para uso no PDF.
Coloque o ficheiro na pasta do projeto, por exemplo:
  Templete - LWSA.pptx
  ou  Template - LWSA.pptx
  ou  assets/Templete - LWSA.pptx
  Plano de fundo fixo (prioritario): Templete - LWSA.png (A4) na pasta do projeto ou assets/
"""
from __future__ import annotations

import glob
import io
import os
import sys
import zipfile
import xml.etree.ElementTree as ET

_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

EMU_PER_MM = 914400 / 25.4


def _tag_local(tag: str) -> str:
    if "}" in tag:
        return tag.rsplit("}", 1)[-1]
    return tag


def _emu_to_mm(emu: int) -> float:
    return float(emu) * 25.4 / 914400.0


def _pptx_presentation_slide_targets(z: zipfile.ZipFile) -> list[str]:
    """Ordem dos slides como no PowerPoint: alvo relativo (ex.: slides/slide8.xml)."""
    pres = ET.fromstring(z.read("ppt/presentation.xml"))
    rels_root = ET.fromstring(z.read("ppt/_rels/presentation.xml.rels"))
    rid_to_target: dict[str, str] = {}
    for rel in rels_root:
        if _tag_local(rel.tag) != "Relationship":
            continue
        rid, tgt = rel.get("Id"), rel.get("Target")
        if rid and tgt:
            rid_to_target[rid] = tgt.replace("\\", "/")
    out: list[str] = []
    for sld in pres.iter():
        if _tag_local(sld.tag) != "sldId":
            continue
        rid = sld.get(f"{{{_NS_R}}}id")
        if not rid:
            continue
        tgt = rid_to_target.get(rid)
        if tgt:
            out.append(tgt)
    return out


def _pptx_slide_part_path(z: zipfile.ZipFile, slide_1based: int) -> str | None:
    if slide_1based < 1:
        return None
    targets = _pptx_presentation_slide_targets(z)
    if slide_1based > len(targets):
        return None
    rel = targets[slide_1based - 1]
    if rel.startswith("/"):
        rel = rel.lstrip("/")
    if not rel.startswith("ppt/"):
        rel = "ppt/" + rel
    return rel if rel in z.namelist() else None


def _slide_size_emu(z: zipfile.ZipFile) -> tuple[int, int]:
    pres = ET.fromstring(z.read("ppt/presentation.xml"))
    for el in pres.iter():
        if _tag_local(el.tag) == "sldSz":
            cx = int(el.get("cx") or 0)
            cy = int(el.get("cy") or 0)
            if cx > 0 and cy > 0:
                return cx, cy
    return 18288000, 10287000


def _grp_ext_cy_off_y(grp: ET.Element) -> tuple[int, int] | None:
    if _tag_local(grp.tag) != "grpSp":
        return None
    grppr = grp.find(f"{{{_NS_P}}}grpSpPr")
    if grppr is None:
        return None
    xfrm = grppr.find(f"{{{_NS_A}}}xfrm")
    if xfrm is None:
        return None
    off = xfrm.find(f"{{{_NS_A}}}off")
    ext = xfrm.find(f"{{{_NS_A}}}ext")
    if off is None or ext is None:
        return None
    try:
        y = int(off.get("y") or 0)
        cy = int(ext.get("cy") or 0)
    except ValueError:
        return None
    if cy <= 0:
        return None
    return cy, y


def _sp_has_srgb_hex(sp: ET.Element, hex_rgb: str) -> bool:
    want = hex_rgb.strip().lstrip("#").upper()
    if len(want) != 6:
        return False
    for srgb in sp.iter():
        if _tag_local(srgb.tag) != "srgbClr":
            continue
        val = (srgb.get("val") or "").strip().lstrip("#").upper()
        if val == want:
            return True
    return False


def _sp_xfrm_ext_bottom(sp: ET.Element) -> tuple[int, int, int] | None:
    """Devolve (cy, y, bottom) em EMU ou None."""
    if _tag_local(sp.tag) != "sp":
        return None
    sppr = sp.find(f"{{{_NS_P}}}spPr")
    if sppr is None:
        return None
    xfrm = sppr.find(f"{{{_NS_A}}}xfrm")
    if xfrm is None:
        return None
    off = xfrm.find(f"{{{_NS_A}}}off")
    ext = xfrm.find(f"{{{_NS_A}}}ext")
    if off is None or ext is None:
        return None
    try:
        y = int(off.get("y") or 0)
        cy = int(ext.get("cy") or 0)
    except ValueError:
        return None
    if cy <= 0:
        return None
    return cy, y, y + cy


def extract_footer_bar_height_mm_from_slide(
    pptx_path: str,
    slide_1based: int = 8,
    footer_srgb_hex: str = "011431",
) -> float | None:
    """
    Altura da faixa de rodapé (mm) a partir do slide indicado do .pptx.
    Procura formas com preenchimento sRGB igual a footer_srgb_hex (#011431 LWSA)
    e usa o `cy` do grupo (grpSp) ancestral que posiciona a faixa no slide —
    o mesmo retângulo visível no template (ex.: slide 8).
    """
    try:
        with zipfile.ZipFile(pptx_path, "r") as z:
            path = _pptx_slide_part_path(z, slide_1based)
            if not path:
                return None
            _, slide_cy = _slide_size_emu(z)
            root = ET.fromstring(z.read(path))
    except (OSError, KeyError, ET.ParseError):
        return None

    tol = max(120000, int(slide_cy * 0.02))
    min_cy = int(0.5 * EMU_PER_MM)
    candidates: list[tuple[int, int]] = []

    def walk(elem: ET.Element, ancestors: list[ET.Element]) -> None:
        chain = ancestors + [elem]
        loc = _tag_local(elem.tag)
        if loc == "sp" and _sp_has_srgb_hex(elem, footer_srgb_hex):
            for g in reversed(chain):
                ge = _grp_ext_cy_off_y(g)
                if ge is None:
                    continue
                cy_g, y_g = ge
                bottom = y_g + cy_g
                if abs(bottom - slide_cy) <= tol and cy_g >= min_cy:
                    candidates.append((cy_g, y_g))
                    break
            else:
                xb = _sp_xfrm_ext_bottom(elem)
                if xb:
                    cy_s, y_s, bot = xb
                    if abs(bot - slide_cy) <= tol and cy_s >= min_cy:
                        candidates.append((cy_s, y_s))
        for child in elem:
            walk(child, chain)

    c_sld = root.find(f"{{{_NS_P}}}cSld")
    start = c_sld.find(f"{{{_NS_P}}}spTree") if c_sld is not None else None
    if start is not None:
        walk(start, [])
    elif c_sld is not None:
        walk(c_sld, [])
    else:
        walk(root, [])

    if not candidates:
        return None
    cy_g, _y_g = max(candidates, key=lambda t: t[1])
    return round(_emu_to_mm(cy_g), 3)


def _pptx_candidates(base_dir: str) -> list[str]:
    names = [
        "Templete - LWSA.pptx",
        "Template - LWSA.pptx",
        os.path.join("assets", "Templete - LWSA.pptx"),
        os.path.join("assets", "Template - LWSA.pptx"),
    ]
    return [os.path.join(base_dir, n) for n in names]


def find_lwsa_pptx(base_dir: str | None = None) -> str | None:
    base_dir = base_dir or os.path.dirname(os.path.abspath(__file__))
    for p in _pptx_candidates(base_dir):
        if os.path.isfile(p):
            return p
    return None


def _png_background_candidates(base_dir: str) -> list[str]:
    names = [
        "Templete - LWSA.png",
        "Template - LWSA.png",
        os.path.join("assets", "Templete - LWSA.png"),
        os.path.join("assets", "Template - LWSA.png"),
    ]
    return [os.path.join(base_dir, n) for n in names]


def find_lwsa_background_png(base_dir: str | None = None) -> str | None:
    """PNG A4 estatico para fundo do PDF (tem prioridade sobre o slide 9 do .pptx)."""
    base_dir = base_dir or os.path.dirname(os.path.abspath(__file__))
    for p in _png_background_candidates(base_dir):
        if os.path.isfile(p):
            return os.path.abspath(p)
    return None


def _srgb_from_element(parent) -> tuple[int, int, int] | None:
    srgb = parent.find(".//a:srgbClr", _NS)
    if srgb is None:
        return None
    val = (srgb.get("val") or "").strip()
    if len(val) != 6:
        return None
    try:
        return tuple(int(val[i : i + 2], 16) for i in (0, 2, 4))
    except ValueError:
        return None


def parse_theme_colors(pptx_path: str) -> dict[str, tuple[int, int, int]]:
    out: dict[str, tuple[int, int, int]] = {}
    with zipfile.ZipFile(pptx_path, "r") as z:
        names = [n for n in z.namelist() if n.startswith("ppt/theme/theme") and n.endswith(".xml")]
        if not names:
            return out
        root = ET.fromstring(z.read(names[0]))
        scheme = root.find(".//a:themeElements/a:clrScheme", _NS)
        if scheme is None:
            return out
        for child in scheme:
            tag = child.tag.split("}")[-1]
            rgb = _srgb_from_element(child)
            if rgb:
                out[tag] = rgb
    return out


def _logo_score(data: bytes) -> float:
    """Prefer imagens horizontais pequenas (típico de logo)."""
    try:
        from PIL import Image
        im = Image.open(io.BytesIO(data))
        w, h = im.size
        if h <= 0:
            return 0.0
        ar = w / float(h)
        if 1.2 <= ar <= 8.0 and h <= 280:
            return ar * 1000 + (300 - h)
        if ar >= 1.0 and h <= 400:
            return ar * 100 + (400 - h)
        return ar * 10
    except Exception:
        return float(len(data)) / 100000.0


WATERMARK_CACHE_NAME = "watermark_slide6.png"

# Plano de fundo A4: slide 9 do template (export COM ou composicao Pillow)
PAGE_BACKGROUND_CACHE_NAME = "slide9_background.png"
SLIDE_PAGE_BACKGROUND = 9


def export_slide_png_powerpoint_com(pptx_path: str, slide_1based: int, out_path: str) -> bool:
    """Exporta um slide para PNG via PowerPoint instalado (Windows). Alta fidelidade."""
    if sys.platform != "win32":
        return False
    try:
        import win32com.client
    except ImportError:
        return False
    pptx_path = os.path.abspath(pptx_path)
    out_path = os.path.abspath(out_path)
    os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
    app = None
    pres = None
    try:
        app = win32com.client.DispatchEx("PowerPoint.Application")
        app.DisplayAlerts = 0
        pres = app.Presentations.Open(pptx_path, WithWindow=True, ReadOnly=True)
        if slide_1based < 1 or slide_1based > pres.Slides.Count:
            return False
        slide = pres.Slides(slide_1based)
        with zipfile.ZipFile(pptx_path, "r") as z:
            sw, sh = _slide_size_emu(z)
        w_px = 2480
        h_px = int(round(w_px * float(sh) / float(sw)))
        slide.Export(out_path, "PNG", w_px, h_px)
        return os.path.isfile(out_path) and os.path.getsize(out_path) > 0
    except Exception:
        return False
    finally:
        try:
            if pres is not None:
                pres.Close()
        except Exception:
            pass
        try:
            if app is not None:
                app.Quit()
        except Exception:
            pass


def render_slide_to_png_pillow(pptx_path: str, slide_1based: int, out_path: str, px_width: int = 2480) -> bool:
    """
    Composicao raster do slide (imagens embutidas + grupos), sem formas vetoriais puras.
    Usado quando o export COM nao esta disponivel.
    """
    try:
        from PIL import Image
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.oxml.ns import qn
    except ImportError:
        return False
    try:
        prs = Presentation(pptx_path)
        slide = prs.slides[slide_1based - 1]
        sw, sh = prs.slide_width, prs.slide_height
        px_h = int(round(px_width * float(sh) / float(sw)))
        canvas = Image.new("RGBA", (px_width, px_h), (255, 255, 255, 255))

        def emu_to_px_x(e: int) -> float:
            return e * px_width / float(sw)

        def emu_to_px_y(e: int) -> float:
            return e * px_h / float(sh)

        def paste_rgba(im: "Image.Image", x: int, y: int) -> None:
            cw, ch = canvas.size
            iw, ih = im.size
            left = max(0, x)
            top = max(0, y)
            right = min(cw, x + iw)
            bottom = min(ch, y + ih)
            if left >= right or top >= bottom:
                return
            crop = im.crop((left - x, top - y, right - x, bottom - y))
            canvas.paste(crop, (left, top), crop)

        def paste_shape_image(shape, ox: int, oy: int) -> None:
            rid = None
            for el in shape._element.iter():
                if el.tag.endswith("}blip"):
                    rid = el.get(qn("r:embed"))
                    if rid:
                        break
            if not rid:
                return
            part = slide.part.related_part(rid)
            im = Image.open(io.BytesIO(part.blob)).convert("RGBA")
            w = max(1, int(emu_to_px_x(int(shape.width))))
            h = max(1, int(emu_to_px_y(int(shape.height))))
            im = im.resize((w, h), Image.Resampling.LANCZOS)
            x = int(emu_to_px_x(ox + int(shape.left)))
            y = int(emu_to_px_y(oy + int(shape.top)))
            paste_rgba(im, x, y)

        def walk(shapes, ox: int, oy: int) -> None:
            for shape in shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    walk(shape.shapes, ox + int(shape.left), oy + int(shape.top))
                else:
                    paste_shape_image(shape, ox, oy)

        walk(slide.shapes, 0, 0)
        os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
        canvas.save(out_path, "PNG")
        return os.path.isfile(out_path) and os.path.getsize(out_path) > 0
    except Exception:
        return False


def extract_page_background_from_slide(
    pptx_path: str,
    cache_dir: str,
    slide_1based: int = SLIDE_PAGE_BACKGROUND,
) -> str | None:
    """
    Gera PNG do slide (tipo A4 no template) para fundo do PDF.

    Por defeito usa composicao Pillow (rapida, sem abrir o PowerPoint).
    Defina LWSA_USE_PPT_COM=1 para exportar via PowerPoint instalado (Windows),
    com maior fidelidade a formas vetoriais (pode abrir o PowerPoint brevemente).
    Se o COM falhar ou nao existir, usa-se Pillow; se Pillow falhar, tenta-se COM.
    """
    dest = os.path.join(cache_dir, PAGE_BACKGROUND_CACHE_NAME)
    use_com_first = os.environ.get("LWSA_USE_PPT_COM", "").lower() in ("1", "true", "yes")
    if use_com_first and export_slide_png_powerpoint_com(pptx_path, slide_1based, dest):
        return dest
    if render_slide_to_png_pillow(pptx_path, slide_1based, dest):
        return dest
    if export_slide_png_powerpoint_com(pptx_path, slide_1based, dest):
        return dest
    return None


def _ppt_norm_media_path(target: str) -> str:
    t = target.replace("\\", "/").lstrip("/")
    if t.startswith("../"):
        t = "ppt/" + t[3:]
    elif not t.startswith("ppt/"):
        t = "ppt/" + t
    return t


def _slide_blip_embed_ids(slide_root: ET.Element) -> list[str]:
    rids: list[str] = []
    for el in slide_root.iter():
        if _tag_local(el.tag) != "blip":
            continue
        rid = el.get(f"{{{_NS_R}}}embed")
        if rid:
            rids.append(rid)
    return rids


def _slide_rels_image_targets(z: zipfile.ZipFile, slide_part: str) -> dict[str, str]:
    """rId -> caminho ppt/media/... para relações de imagem do slide."""
    base = os.path.basename(slide_part)
    rels_path = f"ppt/slides/_rels/{base}.rels"
    if rels_path not in z.namelist():
        return {}
    out: dict[str, str] = {}
    root = ET.fromstring(z.read(rels_path))
    for rel in root:
        if _tag_local(rel.tag) != "Relationship":
            continue
        typ = (rel.get("Type") or "").lower()
        if "image" not in typ:
            continue
        rid, tgt = rel.get("Id"), rel.get("Target")
        if rid and tgt:
            out[rid] = _ppt_norm_media_path(tgt)
    return out


def extract_watermark_from_slide(
    pptx_path: str,
    cache_dir: str,
    slide_1based: int = 6,
    alpha_mult: float = 0.42,
) -> str | None:
    """
    Extrai a primeira imagem referenciada no slide (ex. logo W no slide 6),
    suaviza o canal alpha e grava WATERMARK_CACHE_NAME no cache.
    """
    try:
        with zipfile.ZipFile(pptx_path, "r") as z:
            spath = _pptx_slide_part_path(z, slide_1based)
            if not spath:
                return None
            slide_root = ET.fromstring(z.read(spath))
            rids = _slide_blip_embed_ids(slide_root)
            rid_to_path = _slide_rels_image_targets(z, spath)
            candidates: list[str] = []
            seen: set[str] = set()
            for rid in rids:
                p = rid_to_path.get(rid)
                if p and p in z.namelist() and p not in seen:
                    candidates.append(p)
                    seen.add(p)
            if not candidates:
                for p in rid_to_path.values():
                    if p in z.namelist() and p not in seen:
                        candidates.append(p)
                        seen.add(p)
            if not candidates:
                return None
            media_path = max(candidates, key=lambda p: z.getinfo(p).file_size)
            raw = z.read(media_path)
    except (OSError, KeyError, ET.ParseError):
        return None

    os.makedirs(cache_dir, exist_ok=True)
    dest = os.path.join(cache_dir, WATERMARK_CACHE_NAME)
    try:
        from PIL import Image
        im = Image.open(io.BytesIO(raw)).convert("RGBA")
        if alpha_mult < 0.999:
            alpha = im.split()[3]
            alpha = alpha.point(lambda px: min(255, int(round(px * alpha_mult))))
            im.putalpha(alpha)
        im.save(dest, "PNG")
    except Exception:
        try:
            with open(dest, "wb") as f:
                f.write(raw)
        except OSError:
            return None
    return dest if os.path.isfile(dest) else None


def extract_top_logos(pptx_path: str, cache_dir: str, max_logos: int = 2) -> list[str]:
    """Grava logo_1.ext, logo_2.ext, ... no cache; devolve caminhos por ordem de relevância."""
    exts_ok = ("png", "jpg", "jpeg", "gif")
    scored: list[tuple[float, bytes, str]] = []
    with zipfile.ZipFile(pptx_path, "r") as z:
        for name in z.namelist():
            if not name.startswith("ppt/media/"):
                continue
            low = name.lower()
            if low.endswith(".webp"):
                continue
            if not any(low.endswith("." + e) for e in exts_ok):
                continue
            data = z.read(name)
            scored.append((_logo_score(data), data, name))
    if not scored:
        return []
    scored.sort(key=lambda x: -x[0])
    os.makedirs(cache_dir, exist_ok=True)
    out: list[str] = []
    for idx, (_sc, data, name) in enumerate(scored[:max_logos]):
        ext = name.rsplit(".", 1)[-1].lower()
        if ext == "jpeg":
            ext = "jpg"
        if ext not in exts_ok:
            continue
        dest = os.path.join(cache_dir, f"logo_{idx + 1}.{ext}")
        with open(dest, "wb") as f:
            f.write(data)
        out.append(dest)
    return out


def _mix(a: tuple[int, int, int], b: tuple[int, int, int], t: float) -> tuple[int, int, int]:
    return tuple(int(a[i] * (1 - t) + b[i] * t) for i in range(3))


def _lighten(c: tuple[int, int, int], amount: float) -> tuple[int, int, int]:
    w = (255, 255, 255)
    return _mix(c, w, amount)


def load_branding_for_pdf(base_dir: str | None = None) -> dict | None:
    """
    Devolve dicionário para PDFReport ou None se não houver .pptx nem PNG de fundo.
    Chaves: page_bg, bar, chapter_fill, accent, logo_path (opcional),
    footer_bar_height_mm (opcional, altura da faixa lida no slide 8 do .pptx),
    page_background_path (fundo em todas as páginas),
    skip_footer_bar (opcional, True quando o PNG já inclui faixa de rodapé),
    watermark_path (opcional, slide 6 — só se não houver fundo de página).
    """
    base_dir = base_dir or os.path.dirname(os.path.abspath(__file__))
    png_static = find_lwsa_background_png(base_dir)
    pptx = find_lwsa_pptx(base_dir)
    if not pptx and not png_static:
        return None

    cache_dir = os.path.join(base_dir, "._lwsa_cache")
    os.makedirs(cache_dir, exist_ok=True)
    meta_path = os.path.join(cache_dir, "source_mtime.txt")

    page_bg_path: str | None = None
    skip_footer_bar = False
    if png_static:
        page_bg_path = png_static
        skip_footer_bar = True

    logo_paths: list[str] = []
    watermark_path: str | None = None
    wm_cached = os.path.join(cache_dir, WATERMARK_CACHE_NAME)
    pbg_cached = os.path.join(cache_dir, PAGE_BACKGROUND_CACHE_NAME)

    need_extract = False
    mtime = 0.0
    if pptx:
        mtime = os.path.getmtime(pptx)
        need_extract = True
        if os.path.isfile(meta_path):
            try:
                with open(meta_path, encoding="ascii") as f:
                    old = float(f.read().strip())
                need_extract = old != mtime
            except Exception:
                need_extract = True

        if need_extract:
            for old in glob.glob(os.path.join(cache_dir, "logo_*")):
                try:
                    os.remove(old)
                except OSError:
                    pass
            try:
                if os.path.isfile(wm_cached):
                    os.remove(wm_cached)
            except OSError:
                pass
            try:
                if os.path.isfile(pbg_cached):
                    os.remove(pbg_cached)
            except OSError:
                pass
            logo_paths = extract_top_logos(pptx, cache_dir, max_logos=2)
            if not page_bg_path:
                page_bg_path = extract_page_background_from_slide(
                    pptx, cache_dir, slide_1based=SLIDE_PAGE_BACKGROUND
                )
                if not page_bg_path:
                    watermark_path = extract_watermark_from_slide(pptx, cache_dir, slide_1based=6)
            with open(meta_path, "w", encoding="ascii") as f:
                f.write(str(mtime))
        else:
            logo_paths = []
            for pat in ("logo_1.*", "logo_2.*"):
                hits = sorted(glob.glob(os.path.join(cache_dir, pat)))
                hits = [p for p in hits if p.lower().split(".")[-1] in ("png", "jpg", "jpeg", "gif")]
                if hits:
                    logo_paths.append(hits[0])
            if not page_bg_path:
                if os.path.isfile(pbg_cached):
                    page_bg_path = pbg_cached
                else:
                    page_bg_path = extract_page_background_from_slide(
                        pptx, cache_dir, slide_1based=SLIDE_PAGE_BACKGROUND
                    )
                if not page_bg_path:
                    if os.path.isfile(wm_cached):
                        watermark_path = wm_cached
                    else:
                        watermark_path = extract_watermark_from_slide(pptx, cache_dir, slide_1based=6)

    if pptx:
        colors = parse_theme_colors(pptx)
        accent = colors.get("accent1") or colors.get("accent2") or (0, 82, 155)
        lt1 = colors.get("lt1", (255, 255, 255))
        dk1 = colors.get("dk1", (40, 40, 40))
        page_bg = _mix(lt1, accent, 0.06)
        bar = accent
        chapter_fill = _lighten(accent, 0.72)
        footer_h_mm = extract_footer_bar_height_mm_from_slide(pptx, slide_1based=8)
    else:
        accent = (0, 82, 155)
        lt1 = (255, 255, 255)
        dk1 = (40, 40, 40)
        page_bg = _mix(lt1, accent, 0.06)
        bar = accent
        chapter_fill = _lighten(accent, 0.72)
        footer_h_mm = None

    logo_right = logo_paths[0] if logo_paths else None
    logo_left = logo_paths[1] if len(logo_paths) > 1 else None

    out = {
        "accent": accent,
        "dk1": dk1,
        "lt1": lt1,
        "page_bg": page_bg,
        "bar": bar,
        "chapter_fill": chapter_fill,
        "logo_path": logo_right,
        "logo_path_left": logo_left,
        "source_pptx": pptx,
    }
    if pptx and footer_h_mm is not None:
        out["footer_bar_height_mm"] = footer_h_mm
    if page_bg_path:
        out["page_background_path"] = page_bg_path
    if skip_footer_bar:
        out["skip_footer_bar"] = True
    elif watermark_path:
        out["watermark_path"] = watermark_path
    return out
